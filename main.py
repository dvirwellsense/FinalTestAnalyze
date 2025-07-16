from datetime import datetime

import pandas as pd
import numpy as np
from pathlib import Path
from scipy.ndimage import label, center_of_mass
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
from pptx.dml.color import RGBColor
from collections import Counter, defaultdict
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.chart import XL_LEGEND_POSITION
from math import ceil
from pptx.util import Inches


def extract_second_matrix(file_path):
    df = pd.read_csv(file_path, header=None)
    row_indices = df[df[0].astype(str).str.startswith("row_0")].index.tolist()
    if len(row_indices) < 3:
        raise ValueError(f"Less than 3 matrices found in {file_path}")
    df_matrix = df.iloc[row_indices[1]:row_indices[2]-1, 1:]
    df_matrix = df_matrix.dropna(axis=1, how='all')
    return df_matrix.astype(int).to_numpy()


def extract_third_matrix(file_path):
    df = pd.read_csv(file_path, header=None)
    row_indices = df[df[0].astype(str).str.startswith("row_0")].index.tolist()
    if len(row_indices) < 3:
        raise ValueError(f"Less than 3 matrices found in {file_path}")
    start_idx = row_indices[2]
    df_matrix = df.iloc[start_idx:, 1:]
    df_matrix = df_matrix.dropna(axis=1, how='all')
    return df_matrix.astype(float).to_numpy()


def find_weight_spots(diff_matrix, threshold_ratio=0.1, num_weights=3):
    max_val = np.max(diff_matrix)
    threshold = threshold_ratio * max_val
    binary = (diff_matrix > threshold).astype(int)
    labeled, num_features = label(binary)
    spots = center_of_mass(diff_matrix, labeled, range(1, num_features + 1))
    spot_values = [diff_matrix[int(y)][int(x)] for y, x in spots]
    sorted_spots = sorted(zip(spot_values, spots), reverse=True)[:num_weights]
    return [(int(y), int(x), val) for val, (y, x) in sorted_spots], labeled


def get_stats_for_label(matrix, labeled_matrix, label_id):
    mask = labeled_matrix == label_id
    values = matrix[mask]
    return np.mean(values), np.max(values)


def classify_files_by_avg(folder_path):
    csv_files = list(Path(folder_path).rglob("*.csv"))
    if len(csv_files) != 2:
        raise ValueError(f"Expected exactly 2 CSV files in {folder_path}, found {len(csv_files)}")
    matrices = [extract_third_matrix(f) for f in csv_files]
    avgs = [np.mean(mat[1:-1, 1:]) for mat in matrices]
    if avgs[0] < avgs[1]:
        return csv_files[0], csv_files[1], avgs[0]
    else:
        return csv_files[1], csv_files[0], avgs[1]


def show_debug_plot(mat_diff, spots, folder_name):
    plt.figure(figsize=(8, 6))
    plt.imshow(mat_diff, cmap='hot', interpolation='nearest')
    ys, xs, _ = zip(*spots)
    plt.scatter(xs, ys, marker='x', color='cyan', s=150, label='Detected weights')
    plt.title(f"Weight spots in {folder_name}")
    plt.colorbar(label="Difference")
    plt.legend()
    plt.tight_layout()
    plt.show()


def pressure_to_color(p):
    if p > 75:
        return 'Red'
    elif p > 60:
        return 'Orange'
    elif p > 45:
        return 'Yellow'
    elif p > 30:
        return 'Green'
    elif p > 15:
        return 'Light Blue'
    elif p > 4:
        return 'Dark Blue'
    else:
        return 'None'


def analyze_sheet(folder_path, debug=False):
    empty_file, weight_file, empty_avg_clean = classify_files_by_avg(folder_path)
    mat_empty = extract_third_matrix(empty_file)[1:-1, 1:]
    mat_weight = extract_third_matrix(weight_file)[1:-1, 1:]
    mat_diff = mat_weight - mat_empty

    pressure_matrix = extract_second_matrix(weight_file)[1:-1, 1:]

    mean_diff = np.mean(mat_diff)
    std_diff = np.std(mat_diff)
    total_diff = np.sum(mat_diff)

    spots, labeled = find_weight_spots(mat_diff)
    sorted_spots = sorted(spots, key=lambda p: p[0])  # sort by row (y)

    weights_lbs = [5, 10, 20]
    weight_responses = {}

    for w, (y, x, _) in zip(weights_lbs, sorted_spots):
        label_id = labeled[y, x]
        # Difference response
        mean_diff_val, max_diff_val = get_stats_for_label(mat_diff, labeled, label_id)
        weight_responses[f"{w}lb_response"] = mean_diff_val
        weight_responses[f"{w}lb_max"] = max_diff_val
        # Absolute values on weight matrix
        mean_abs_val, max_abs_val = get_stats_for_label(mat_weight, labeled, label_id)
        weight_responses[f"{w}lb_avg_val"] = mean_abs_val
        weight_responses[f"{w}lb_max_val"] = max_abs_val

        mean_pressure, max_pressure = get_stats_for_label(pressure_matrix, labeled, label_id)
        weight_responses[f"{w}lb_color"] = pressure_to_color(max_pressure)

    if debug:
        print(f"Debug info for {Path(folder_path).name}:")
        for w in weights_lbs:
            print(f"  {w}lb response: mean={weight_responses[f'{w}lb_response']:.4e}, max={weight_responses[f'{w}lb_max']:.4e}")
            print(f"  {w}lb value: mean={weight_responses[f'{w}lb_avg_val']:.4e}, max={weight_responses[f'{w}lb_max_val']:.4e}")
        show_debug_plot(mat_diff, spots, Path(folder_path).name)

    return {
        "Mat": Path(folder_path).name,
        "mean_diff": mean_diff,
        "std_diff": std_diff,
        "total_diff": total_diff,
        "empty_avg_clean": empty_avg_clean,
        **weight_responses,
        "ratio_10_to_20": weight_responses["10lb_response"] / weight_responses["20lb_response"],
        "ratio_5_to_20": weight_responses["5lb_response"] / weight_responses["20lb_response"],
        "weight_file": weight_file.name
    }


def add_mat_slide(prs, row, base_folder):
        layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(layout)

        # מילון צבעים
        COLOR_MAP = {
            'Red': RGBColor(255, 0, 0),
            'Orange': RGBColor(255, 165, 0),
            'Yellow': RGBColor(255, 255, 0),
            'Green': RGBColor(0, 128, 0),
            'Light Blue': RGBColor(173, 216, 230),
            'Dark Blue': RGBColor(0, 0, 139),
            'None': RGBColor(255, 255, 255)
        }

        def color_cell(cell, color_name):
            # צבע רקע
            cell.fill.solid()
            cell.fill.fore_color.rgb = COLOR_MAP[color_name]
            # צבע טקסט אם כהה
            if color_name in {'Red', 'Dark Blue'}:
                for paragraph in cell.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.color.rgb = RGBColor(255, 255, 255)

        # טבלה 1 – תגובות הפרש
        table1 = slide.shapes.add_table(6, 2, Inches(0.5), Inches(0.8), Inches(4.5), Inches(2.2)).table
        table1.columns[0].width = Inches(2)
        table1.columns[1].width = Inches(2.5)
        table1.cell(0, 0).text = "Baseline avg (no weight)"
        table1.cell(0, 1).text = f"{row['empty_avg_clean']:.2e}"
        table1.cell(1, 0).text = "5lb: avr / max [response]"
        table1.cell(1, 1).text = f"{row['5lb_response']:.2e} / {row['5lb_max']:.2e}"
        table1.cell(2, 0).text = "10lb: avr / max [response]"
        table1.cell(2, 1).text = f"{row['10lb_response']:.2e} / {row['10lb_max']:.2e}"
        table1.cell(3, 0).text = "20lb: avr / max [response]"
        table1.cell(3, 1).text = f"{row['20lb_response']:.2e} / {row['20lb_max']:.2e}"
        table1.cell(4, 0).text = "Ratio 10lb / 20lb"
        table1.cell(4, 1).text = f"{row['ratio_10_to_20']:.2f}"
        table1.cell(5, 0).text = "Ratio 5lb / 20lb"
        table1.cell(5, 1).text = f"{row['ratio_5_to_20']:.2f}"

        # טבלה 2 – ערכים מוחלטים + צבע עם שורה עליונה של כותרות
        table2 = slide.shapes.add_table(4, 3, Inches(0.5), Inches(5.1), Inches(4.5), Inches(1.3)).table
        table2.columns[0].width = Inches(2)
        table2.columns[1].width = Inches(2.5)
        table2.columns[2].width = Inches(1.5)

        # כותרות
        table2.cell(0, 0).text = "Weight"
        table2.cell(0, 1).text = "avg / max"
        table2.cell(0, 2).text = "Max color"

        # שורה 1
        table2.cell(1, 0).text = "5lb"
        table2.cell(1, 1).text = f"{row['5lb_avg_val']:.2e} / {row['5lb_max_val']:.2e}"
        table2.cell(1, 2).text = row["5lb_color"]
        color_cell(table2.cell(1, 2), row["5lb_color"])

        # שורה 2
        table2.cell(2, 0).text = "10lb"
        table2.cell(2, 1).text = f"{row['10lb_avg_val']:.2e} / {row['10lb_max_val']:.2e}"
        table2.cell(2, 2).text = row["10lb_color"]
        color_cell(table2.cell(2, 2), row["10lb_color"])

        # שורה 3
        table2.cell(3, 0).text = "20lb"
        table2.cell(3, 1).text = f"{row['20lb_avg_val']:.2e} / {row['20lb_max_val']:.2e}"
        table2.cell(3, 2).text = row["20lb_color"]
        color_cell(table2.cell(3, 2), row["20lb_color"])

        # תמונה
        folder = Path(base_folder) / row["Mat"]
        image_file = folder / row["weight_file"].replace("_rawData.csv", "_heatmap.png")
        if image_file.exists():
            pic = slide.shapes.add_picture(str(image_file), Inches(6.75), Inches(0.5), width=Inches(3))
            pic.top = prs.slide_height - pic.height

        # כותרת
        textbox = slide.shapes.add_textbox(Inches(0.5), Inches(0.1), Inches(8), Inches(0.2))
        tf = textbox.text_frame
        p = tf.add_paragraph()
        p.text = f"Mat {row['Mat']}"
        p.font.size = Pt(24)
        p.font.bold = True
        p.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER


def add_summary_slide(prs, df):
    max_rows_per_slide = 14
    total_rows = len(df)
    num_slides = ceil(total_rows / max_rows_per_slide)

    headers = ["Mat", "20lb_avr", "10lb_avr", "5lb_avr", "20lb_resp", "10lb_resp", "5lb_resp", "Baseline"]

    for slide_idx in range(num_slides):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        textbox = slide.shapes.add_textbox(Inches(0.5), Inches(0.1), Inches(8), Inches(0.4))
        tf = textbox.text_frame
        p = tf.add_paragraph()
        p.text = f"Summary Table (Page {slide_idx + 1})"
        p.font.size = Pt(24)
        p.font.bold = True
        p.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER

        # חיתוך ה-DataFrame לחלק המתאים
        start_idx = slide_idx * max_rows_per_slide
        end_idx = min(start_idx + max_rows_per_slide, total_rows)
        df_slice = df.iloc[start_idx:end_idx]

        rows = len(df_slice) + 1  # +1 לשורה של הכותרות
        cols = len(headers)

        table = slide.shapes.add_table(rows, cols, Inches(0.2), Inches(1), Inches(9), Inches(0.3 * rows)).table

        # הוספת כותרות
        for col, h in enumerate(headers):
            table.cell(0, col).text = h

        # תוכן הטבלה
        for i, (_, row) in enumerate(df_slice.iterrows(), start=1):
            table.cell(i, 0).text = str(row["Mat"])
            table.cell(i, 1).text = f"{row['20lb_avg_val']:.2e}"
            table.cell(i, 2).text = f"{row['10lb_avg_val']:.2e}"
            table.cell(i, 3).text = f"{row['5lb_avg_val']:.2e}"
            table.cell(i, 4).text = f"{row['20lb_response']:.2e}"
            table.cell(i, 5).text = f"{row['10lb_response']:.2e}"
            table.cell(i, 6).text = f"{row['5lb_response']:.2e}"
            table.cell(i, 7).text = f"{row['empty_avg_clean']:.2e}"


def add_average_value_chart_slide(prs, df):
    # ממיין את הדאטה לפי מספר מזרן כערך מספרי
    df_sorted = df.copy()
    df_sorted["MatNum"] = df_sorted["Mat"].astype(int)
    df_sorted = df_sorted.sort_values("MatNum")

    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # כותרת השקופית
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.1), Inches(8), Inches(0.4))
    tf = title_box.text_frame
    p = tf.add_paragraph()
    p.text = "Average Capacitance Trend by Mat Number"
    p.font.size = Pt(24)
    p.font.bold = True
    p.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER

    # הכנת הנתונים לגרף קווי
    chart_data = CategoryChartData()
    chart_data.categories = list(df_sorted["MatNum"])
    chart_data.add_series("20lb avg", list(df_sorted["20lb_avg_val"]))
    chart_data.add_series("10lb avg", list(df_sorted["10lb_avg_val"]))
    chart_data.add_series("5lb avg", list(df_sorted["5lb_avg_val"]))

    x, y = Inches(0.5), Inches(1.0)
    width, height = Inches(9), Inches(4.5)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.LINE_MARKERS,
        x, y, width, height,
        chart_data
    ).chart

    # עיצוב
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.chart_title.text_frame.text = "Average Capacitance by Mat Number"
    chart.value_axis.has_major_gridlines = True


def add_color_summary_slide(prs, color_counters):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    textbox = slide.shapes.add_textbox(Inches(0.5), Inches(0.1), Inches(8), Inches(0.4))
    tf = textbox.text_frame
    p = tf.add_paragraph()
    p.text = "Color Distribution by Weight"
    p.font.size = Pt(24)
    p.font.bold = True
    p.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER

    # צבעים בפועל
    color_map = {
        "Red": RGBColor(255, 0, 0),
        "Orange": RGBColor(255, 165, 0),
        "Yellow": RGBColor(255, 255, 0),
        "Green": RGBColor(0, 128, 0),
        "Light Blue": RGBColor(173, 216, 230),
        "Dark Blue": RGBColor(0, 0, 139),
        "White": RGBColor(255, 255, 255),
        "None": RGBColor(200, 200, 200)
    }

    # הצבעים המצופים לפי משקל
    expected_colors = {
        5: "Light Blue",
        10: "Green",
        20: "Orange"
    }

    left_positions = [Inches(0.5), Inches(3.5), Inches(6.5)]
    chart_width = Inches(3.15)
    chart_height = Inches(3)

    for i, weight in enumerate([5, 10, 20]):
        counts = color_counters[f"{weight}lb"]
        if not counts:
            continue

        chart_data = CategoryChartData()
        labels = list(counts.keys())
        values = list(counts.values())
        chart_data.categories = labels
        chart_data.add_series(f"{weight}lb", values)

        x, y = left_positions[i], Inches(1.2)
        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.PIE,
            x, y, chart_width, chart_height,
            chart_data
        ).chart

        chart.has_data_labels = True
        data_labels = chart.plots[0].series[0].data_labels
        data_labels.show_percentage = True
        chart.has_legend = True
        chart.legend.include_in_layout = False
        chart.legend.position = XL_LEGEND_POSITION.RIGHT

        # ✅ הוספת הצבע המצופה לכותרת
        expected = expected_colors[weight]
        chart.has_title = True
        chart.chart_title.text_frame.clear()
        p = chart.chart_title.text_frame.paragraphs[0]
        p.text = f"{weight}lb\n(Expected: {expected})"
        p.font.size = Pt(18)

        # 🎨 צביעת פרוסות הפאי
        for j, point in enumerate(chart.plots[0].series[0].points):
            label = labels[j]
            rgb = color_map.get(label, RGBColor(128, 128, 128))  # ברירת מחדל אפור
            point.format.fill.solid()
            point.format.fill.fore_color.rgb = rgb

        thresholds = [
            ("White", "0"),
            ("Dark Blue", "4"),
            ("Light Blue", "15"),
            ("Green", "30"),
            ("Yellow", "45"),
            ("Orange", "60"),
            ("Red", "75")
        ]

        # הוספת הטבלה (7 שורות, 2 עמודות)
        # יצירת הטבלה במקום מבוקש בגודל מתאים
        table_shape = slide.shapes.add_table(rows=7, cols=2, left=Inches(0.5), top=Inches(4.5), width=Inches(2.4),
                                             height=Inches(2.0))
        table = table_shape.table

        # הגדרת רוחב עמודות
        table.columns[0].width = Inches(1.4)
        table.columns[1].width = Inches(1.2)

        # מילוי טבלה
        for i, (color_name, value) in enumerate(thresholds):
            cell = table.cell(i, 0)
            cell.text = color_name
            cell.fill.solid()
            cell.fill.fore_color.rgb = color_map[color_name]

            # להפוך את הטקסט בלבן אם הרקע כהה
            if color_name in {"Red", "Dark Blue", "Green"}:
                for paragraph in cell.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.color.rgb = RGBColor(255, 255, 255)
            else:
                for paragraph in cell.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.color.rgb = RGBColor(0, 0, 0)

            table.cell(i, 1).text = value

    # יצירת הטבלה החדשה (4 שורות, 4 עמודות)
    table_shape2 = slide.shapes.add_table(rows=4, cols=4, left=Inches(3.2), top=Inches(5.7), width=Inches(5.2),
                                          height=Inches(1.5))
    table2 = table_shape2.table

    # כותרות
    headers = ["libra", "kg", "cm^2", "mmHg"]
    for col, title in enumerate(headers):
        table2.cell(0, col).text = title

    # נתונים
    data = [
        [5, "2.268", "100", 16.682],
        [10, "4.5359", "100", 33.364],
        [20, "9.0718", "100", 66.729]
    ]

    # מילוי נתונים
    for row_idx, row_data in enumerate(data, start=1):
        for col_idx, value in enumerate(row_data):
            cell = table2.cell(row_idx, col_idx)
            cell.text = str(value)
            # צבע לפי mmHg (עמודה אחרונה בלבד)
            if col_idx == 3:
                color_name = pressure_to_color(float(value))
                cell.fill.solid()
                cell.fill.fore_color.rgb = color_map[color_name]
                # טקסט לבן אם כהה
                if color_name in {"Red", "Blue"}:
                    for paragraph in cell.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.color.rgb = RGBColor(255, 255, 255)


def create_ppt_report(df, base_folder, color_counters):
    prs = Presentation()
    for _, row in df.iterrows():
        add_mat_slide(prs, row, base_folder)
    add_summary_slide(prs, df)
    add_average_value_chart_slide(prs, df)
    add_color_summary_slide(prs, color_counters)
    # קבלת התאריך הנוכחי במחרוזת
    date_str = datetime.now().strftime("%Y-%m-%d_%H-%M-%S")

    # יצירת שם קובץ עם התאריך
    output_path = Path(base_folder) / f"final_test_analysis_report_{date_str}.pptx"

    prs.save(output_path)
    print(f"PowerPoint report saved to {output_path}")


def analyze_all(base_folder, debug=False):
    results = []
    color_counters = defaultdict(Counter)
    for subfolder in Path(base_folder).iterdir():
        if subfolder.is_dir():
            try:
                print(f"Analyzing {subfolder.name}...")
                result = analyze_sheet(subfolder, debug=debug)
                results.append(result)
                for w in [5, 10, 20]:
                    color = result[f"{w}lb_color"]
                    color_counters[f"{w}lb"][color] += 1
            except Exception as e:
                print(f"Error in {subfolder.name}: {e}")

    df = pd.DataFrame(results)
    create_ppt_report(df, base_folder, color_counters)


if __name__ == "__main__":
    base_folder = r"C:\Users\dvirs\Desktop\wellsense-VU-Player\RawData\N\A\new_label"
    analyze_all(base_folder, False)
