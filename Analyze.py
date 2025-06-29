import pandas as pd
import numpy as np
from pathlib import Path
from scipy.ndimage import label, center_of_mass
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import io

def extract_third_matrix(file_path):
    df = pd.read_csv(file_path, header=None)
    row_indices = df[df[0].astype(str).str.startswith("row_0")].index.tolist()
    if len(row_indices) < 3:
        raise ValueError(f"Less than 3 matrices found in {file_path}")
    start_idx = row_indices[2]
    df_matrix = df.iloc[start_idx:, 1:]
    df_matrix = df_matrix.dropna(axis=1, how='all')
    return df_matrix.astype(float).to_numpy()

def find_weight_spots(diff_matrix, threshold_ratio=0.1, num_weights=3):
    max_val = np.max(diff_matrix)
    threshold = threshold_ratio * max_val
    binary = (diff_matrix > threshold).astype(int)
    labeled, num_features = label(binary)
    spots = center_of_mass(diff_matrix, labeled, range(1, num_features + 1))
    spot_values = [diff_matrix[int(y)][int(x)] for y, x in spots]
    sorted_spots = sorted(zip(spot_values, spots), reverse=True)[:num_weights]
    return [(int(y), int(x), val) for val, (y, x) in sorted_spots], labeled

def get_stats_for_label(matrix, labeled_matrix, label_id):
    mask = labeled_matrix == label_id
    values = matrix[mask]
    return np.mean(values), np.max(values)

def classify_files_by_avg(folder_path):
    csv_files = list(Path(folder_path).rglob("*.csv"))
    if len(csv_files) != 2:
        raise ValueError(f"Expected exactly 2 CSV files in {folder_path}, found {len(csv_files)}")
    matrices = [extract_third_matrix(f) for f in csv_files]
    avgs = [np.mean(mat[1:-1, 1:]) for mat in matrices]
    if avgs[0] < avgs[1]:
        return csv_files[0], csv_files[1], avgs[0]
    else:
        return csv_files[1], csv_files[0], avgs[1]

def show_debug_plot(mat_diff, spots, folder_name):
    plt.figure(figsize=(8, 6))
    plt.imshow(mat_diff, cmap='hot', interpolation='nearest')
    ys, xs, _ = zip(*spots)
    plt.scatter(xs, ys, marker='x', color='cyan', s=150, label='Detected weights')
    plt.title(f"Weight spots in {folder_name}")
    plt.colorbar(label="Difference")
    plt.legend()
    plt.tight_layout()
    plt.show()

def analyze_sheet(folder_path, debug=False):
    empty_file, weight_file, empty_avg_clean = classify_files_by_avg(folder_path)
    mat_empty = extract_third_matrix(empty_file)[1:-1, 1:]
    mat_weight = extract_third_matrix(weight_file)[1:-1, 1:]
    mat_diff = mat_weight - mat_empty

    mean_diff = np.mean(mat_diff)
    std_diff = np.std(mat_diff)
    total_diff = np.sum(mat_diff)

    spots, labeled = find_weight_spots(mat_diff)
    sorted_spots = sorted(spots, key=lambda p: p[0])  # sort by row (y)

    weights_lbs = [20, 10, 5]
    weight_responses = {}

    for w, (y, x, _) in zip(weights_lbs, sorted_spots):
        label_id = labeled[y, x]
        # Difference response
        mean_diff_val, max_diff_val = get_stats_for_label(mat_diff, labeled, label_id)
        weight_responses[f"{w}lb_response"] = mean_diff_val
        weight_responses[f"{w}lb_max"] = max_diff_val
        # Absolute values on weight matrix
        mean_abs_val, max_abs_val = get_stats_for_label(mat_weight, labeled, label_id)
        weight_responses[f"{w}lb_avg_val"] = mean_abs_val
        weight_responses[f"{w}lb_max_val"] = max_abs_val

    if debug:
        print(f"Debug info for {Path(folder_path).name}:")
        for w in weights_lbs:
            print(f"  {w}lb response: mean={weight_responses[f'{w}lb_response']:.4e}, max={weight_responses[f'{w}lb_max']:.4e}")
            print(f"  {w}lb value: mean={weight_responses[f'{w}lb_avg_val']:.4e}, max={weight_responses[f'{w}lb_max_val']:.4e}")
        show_debug_plot(mat_diff, spots, Path(folder_path).name)

    return {
        "Mat": Path(folder_path).name,
        "mean_diff": mean_diff,
        "std_diff": std_diff,
        "total_diff": total_diff,
        "empty_avg_clean": empty_avg_clean,
        **weight_responses,
        "ratio_10_to_20": weight_responses["10lb_response"] / weight_responses["20lb_response"],
        "ratio_5_to_20": weight_responses["5lb_response"] / weight_responses["20lb_response"],
        "weight_file": weight_file.name
    }

def add_mat_slide(prs, row, base_folder):
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)

    # Table 1: Differences
    table1 = slide.shapes.add_table(6, 2, Inches(0.5), Inches(0.8), Inches(4.5), Inches(2.2)).table
    table1.columns[0].width = Inches(2)
    table1.columns[1].width = Inches(2.5)
    table1.cell(0, 0).text = "Baseline avg (no weight)"
    table1.cell(0, 1).text = f"{row['empty_avg_clean']:.2e}"
    table1.cell(1, 0).text = "20lb: avr / max [response]"
    table1.cell(1, 1).text = f"{row['20lb_response']:.2e} / {row['20lb_max']:.2e}"
    table1.cell(2, 0).text = "10lb: avr / max [response]"
    table1.cell(2, 1).text = f"{row['10lb_response']:.2e} / {row['10lb_max']:.2e}"
    table1.cell(3, 0).text = "5lb: avr / max [response]"
    table1.cell(3, 1).text = f"{row['5lb_response']:.2e} / {row['5lb_max']:.2e}"
    table1.cell(4, 0).text = "Ratio 10lb / 20lb"
    table1.cell(4, 1).text = f"{row['ratio_10_to_20']:.2f}"
    table1.cell(5, 0).text = "Ratio 5lb / 20lb"
    table1.cell(5, 1).text = f"{row['ratio_5_to_20']:.2f}"

    # Table 2: Absolute values
    table2 = slide.shapes.add_table(3, 2, Inches(0.5), Inches(5.1), Inches(4.5), Inches(1.0)).table
    table2.columns[0].width = Inches(2)
    table2.columns[1].width = Inches(2.5)
    table2.cell(0, 0).text = "20lb value (avg / max)"
    table2.cell(0, 1).text = f"{row['20lb_avg_val']:.2e} / {row['20lb_max_val']:.2e}"
    table2.cell(1, 0).text = "10lb value (avg / max)"
    table2.cell(1, 1).text = f"{row['10lb_avg_val']:.2e} / {row['10lb_max_val']:.2e}"
    table2.cell(2, 0).text = "5lb value (avg / max)"
    table2.cell(2, 1).text = f"{row['5lb_avg_val']:.2e} / {row['5lb_max_val']:.2e}"

    # Image
    folder = Path(base_folder) / row["Mat"]
    image_file = folder / row["weight_file"].replace("_rawData.csv", "_heatmap.png")
    if image_file.exists():
        pic = slide.shapes.add_picture(str(image_file), Inches(6), Inches(0.5), width=Inches(3))
        pic.top = prs.slide_height - pic.height

    # Title
    textbox = slide.shapes.add_textbox(Inches(0.5), Inches(0.1), Inches(8), Inches(0.2))
    tf = textbox.text_frame
    p = tf.add_paragraph()
    p.text = f"Mat {row['Mat']}"
    p.font.size = Pt(24)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

def add_summary_slide(prs, df):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Summary Table"
    rows, cols = len(df) + 1, 7
    table = slide.shapes.add_table(rows, cols, Inches(0.2), Inches(1), Inches(9), Inches(0.3 * rows)).table
    headers = ["Mat", "20lb", "10lb", "5lb", "10/20", "5/20", "Baseline"]
    for col, h in enumerate(headers):
        table.cell(0, col).text = h
    for i, (_, row) in enumerate(df.iterrows(), start=1):
        table.cell(i, 0).text = str(row["Mat"])
        table.cell(i, 1).text = f"{row['20lb_response']:.2e}"
        table.cell(i, 2).text = f"{row['10lb_response']:.2e}"
        table.cell(i, 3).text = f"{row['5lb_response']:.2e}"
        table.cell(i, 4).text = f"{row['ratio_10_to_20']:.2f}"
        table.cell(i, 5).text = f"{row['ratio_5_to_20']:.2f}"
        table.cell(i, 6).text = f"{row['empty_avg_clean']:.2e}"

def save_excel_with_summary(df, base_folder):
    with pd.ExcelWriter(base_folder, engine='xlsxwriter') as writer:
        # Write full data
        df.to_excel(writer, sheet_name='Summary', index=False)
        workbook  = writer.book
        worksheet = writer.sheets['Summary']

        # Set column widths for readability
        for i, col in enumerate(df.columns):
            column_width = max(len(str(col)), 15)
            worksheet.set_column(i, i, column_width)

        # Create chart for response values
        chart_resp = workbook.add_chart({'type': 'line'})
        chart_abs = workbook.add_chart({'type': 'line'})

        def add_series(chart, name, col_index, color):
            chart.add_series({
                'name':       ["Summary", 0, col_index],
                'categories': ["Summary", 1, 0, len(df), 0],
                'values':     ["Summary", 1, col_index, len(df), col_index],
                'line':       {'color': color},
            })

        # Define color palette
        colors = {
            "20lb_response": 'red',
            "10lb_response": 'blue',
            "5lb_response": 'green',
            "20lb_avg_val": 'darkred',
            "10lb_avg_val": 'darkblue',
            "5lb_avg_val": 'darkgreen'
        }

        # Add response series
        for col in ["20lb_response", "10lb_response", "5lb_response"]:
            add_series(chart_resp, col, df.columns.get_loc(col), colors[col])

        chart_resp.set_title({'name': 'Weight Responses'})
        chart_resp.set_x_axis({'name': 'Mat'})
        chart_resp.set_y_axis({'name': 'Mean Response'})
        worksheet.insert_chart('J2', chart_resp, {'x_scale': 1.5, 'y_scale': 1.5})

        # Add absolute value series
        for col in ["20lb_avg_val", "10lb_avg_val", "5lb_avg_val"]:
            add_series(chart_abs, col, df.columns.get_loc(col), colors[col])

        chart_abs.set_title({'name': 'Absolute Capacitance Values'})
        chart_abs.set_x_axis({'name': 'Mat'})
        chart_abs.set_y_axis({'name': 'Mean Capacitance'})
        worksheet.insert_chart('J20', chart_abs, {'x_scale': 1.5, 'y_scale': 1.5})

def create_ppt_report(df, base_folder):
    prs = Presentation()
    for _, row in df.iterrows():
        add_mat_slide(prs, row, base_folder)
    add_summary_slide(prs, df)
    output_path = Path(base_folder) / "final_test_analysis_report.pptx"
    prs.save(output_path)
    print(f"PowerPoint report saved to {output_path}")

def analyze_all(base_folder, debug=False):
    results = []
    for subfolder in Path(base_folder).iterdir():
        if subfolder.is_dir():
            try:
                print(f"Analyzing {subfolder.name}...")
                result = analyze_sheet(subfolder, debug=debug)
                results.append(result)
            except Exception as e:
                print(f"Error in {subfolder.name}: {e}")

    df = pd.DataFrame(results)
    create_ppt_report(df, base_folder)
    xlsx_path = Path(base_folder) / "final_test_analysis_summary.xlsx"
    save_excel_with_summary(df, xlsx_path)
    print(f"Excel summary saved to {xlsx_path}")

if __name__ == "__main__":
    import argparse
    parser = argparse.ArgumentParser(description="Analyze final test data from calibration mats.")
    parser.add_argument("folder", help="Path to folder containing mats subfolders")
    parser.add_argument("--debug", action="store_true", help="Enable debug output and visualization")
    args = parser.parse_args()
    analyze_all(args.folder, debug=args.debug)
